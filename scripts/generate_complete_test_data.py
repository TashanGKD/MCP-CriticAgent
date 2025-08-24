#!/usr/bin/env python3
"""
完整测试数据生成器 - 包含所有格式

生成用于MCP工具测试的完整测试数据集：
- CSV: 多种业务数据
- Excel: 多工作表复杂数据
- JSON: 配置和结果数据
- PDF: 报告文档
- Word: 富文本文档
- PowerPoint: 演示文稿
- 图片: 测试图像
- XML: 结构化数据

作者: AI Assistant  
日期: 2025-08-24
"""

import pandas as pd
from pathlib import Path
from datetime import datetime, timedelta
import random
import json

# 高级格式库
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches
from PIL import Image, ImageDraw, ImageFont
import xml.etree.ElementTree as ET

# 测试数据目录
TEST_DATA_DIR = Path(__file__).parent.parent / "test_data"
TEST_DATA_DIR.mkdir(exist_ok=True)

def generate_comprehensive_excel():
    """生成全面的Excel测试文件"""
    print("📊 生成综合Excel文件...")
    
    # 1. 财务报表Excel
    with pd.ExcelWriter(TEST_DATA_DIR / "financial_report.xlsx", engine='openpyxl') as writer:
        # 月度收入数据
        months = ['1月', '2月', '3月', '4月', '5月', '6月', '7月', '8月', '9月', '10月', '11月', '12月']
        revenue_data = pd.DataFrame({
            '月份': months,
            '产品A收入': [random.randint(45000, 65000) for _ in months],
            '产品B收入': [random.randint(25000, 45000) for _ in months], 
            '产品C收入': [random.randint(15000, 35000) for _ in months],
            '服务收入': [random.randint(10000, 25000) for _ in months]
        })
        revenue_data['总收入'] = revenue_data[['产品A收入', '产品B收入', '产品C收入', '服务收入']].sum(axis=1)
        revenue_data.to_excel(writer, sheet_name='月度收入', index=False)
        
        # 成本明细
        cost_data = pd.DataFrame({
            '月份': months,
            '人工成本': [random.randint(20000, 30000) for _ in months],
            '材料成本': [random.randint(15000, 25000) for _ in months],
            '运营成本': [random.randint(8000, 15000) for _ in months],
            '营销成本': [random.randint(5000, 12000) for _ in months],
            '其他成本': [random.randint(3000, 8000) for _ in months]
        })
        cost_data['总成本'] = cost_data[['人工成本', '材料成本', '运营成本', '营销成本', '其他成本']].sum(axis=1)
        cost_data.to_excel(writer, sheet_name='月度成本', index=False)
        
        # 利润分析
        profit_data = pd.DataFrame({
            '月份': months,
            '总收入': revenue_data['总收入'],
            '总成本': cost_data['总成本'],
        })
        profit_data['毛利润'] = profit_data['总收入'] - profit_data['总成本']
        profit_data['利润率'] = (profit_data['毛利润'] / profit_data['总收入'] * 100).round(2)
        profit_data.to_excel(writer, sheet_name='利润分析', index=False)
    
    # 2. 员工管理Excel
    departments = ['技术部', '产品部', '设计部', '运营部', '市场部', '人事部', '财务部']
    positions = ['工程师', '高级工程师', '技术专家', '经理', '总监', '专员', '主管']
    employee_data = []
    
    for i in range(50):
        employee_data.append({
            'ID': f'EMP{i+1:03d}',
            '姓名': f'员工{i+1}',
            '部门': random.choice(departments),
            '职位': random.choice(positions),
            '薪资': random.randint(8000, 35000),
            '入职日期': (datetime(2020, 1, 1) + timedelta(days=random.randint(0, 1500))).strftime('%Y-%m-%d'),
            '年龄': random.randint(22, 45),
            '学历': random.choice(['本科', '硕士', '博士', '专科']),
            '工作年限': random.randint(0, 20),
            '绩效评级': random.choice(['A', 'B', 'C', 'D'])
        })
    
    employee_df = pd.DataFrame(employee_data)
    employee_df.to_excel(TEST_DATA_DIR / "employee_management.xlsx", index=False)
    
    print("✅ Excel文件生成完成")

def generate_comprehensive_pdf():
    """生成PDF报告文件"""
    print("📄 生成PDF报告...")
    
    # 创建PDF报告
    c = canvas.Canvas(str(TEST_DATA_DIR / "mcp_test_report.pdf"), pagesize=A4)
    width, height = A4
    
    # 标题页
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width/2, height - 100, "MCP测试框架报告")
    
    c.setFont("Helvetica", 16)
    c.drawCentredString(width/2, height - 150, "Model Context Protocol工具测试")
    c.drawCentredString(width/2, height - 180, f"生成日期: {datetime.now().strftime('%Y年%m月%d日')}")
    
    # 添加内容
    y_position = height - 250
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, y_position, "测试概览")
    
    y_position -= 30
    c.setFont("Helvetica", 12)
    content = [
        "• 测试工具总数: 25个",
        "• 成功测试: 20个 (80%)",
        "• 失败测试: 5个 (20%)",
        "• 平均测试时间: 15.3秒",
        "• 支持运行时: npx, uvx",
        "",
        "主要发现:",
        "• uvx类型工具部署成功率达到90%",
        "• Excel处理类工具需要绝对路径",
        "• AI智能测试生成准确率为85%",
        "• 大部分工具支持stdio传输模式"
    ]
    
    for line in content:
        c.drawString(70, y_position, line)
        y_position -= 20
    
    # 添加表格数据
    y_position -= 40
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y_position, "测试结果详情")
    
    y_position -= 30
    c.setFont("Helvetica", 10)
    table_data = [
        ["工具名称", "运行时", "状态", "耗时(s)", "工具数"],
        ["Excel MCP Server", "uvx", "成功", "12.5", "25"],
        ["Context7", "npx", "成功", "8.3", "2"],
        ["Blender MCP", "uvx", "成功", "18.7", "8"],
        ["ElevenLabs MCP", "uvx", "失败", "30.0", "0"]
    ]
    
    for row in table_data:
        x_offset = 50
        for cell in row:
            c.drawString(x_offset, y_position, str(cell))
            x_offset += 120
        y_position -= 15
    
    c.save()
    print("✅ PDF报告生成完成")

def generate_word_document():
    """生成Word文档"""
    print("📝 生成Word文档...")
    
    doc = Document()
    
    # 添加标题
    title = doc.add_heading('MCP工具测试指南', 0)
    
    # 添加介绍段落
    intro = doc.add_paragraph('本文档介绍如何使用MCP测试框架对Model Context Protocol工具进行自动化测试。')
    
    # 添加子标题
    doc.add_heading('1. 快速开始', level=1)
    
    # 添加步骤列表
    doc.add_paragraph('安装依赖:', style='List Number')
    doc.add_paragraph('uv sync', style='List Bullet')
    
    doc.add_paragraph('运行测试:', style='List Number') 
    doc.add_paragraph('uv run python -m src.main test-url "https://github.com/example/tool"', style='List Bullet')
    
    # 添加表格
    doc.add_heading('2. 支持的运行时', level=1)
    
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Light Shading Accent 1'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '运行时'
    hdr_cells[1].text = '描述'
    hdr_cells[2].text = '示例命令'
    
    runtime_data = [
        ('npx', 'Node.js包执行器', 'npx @example/tool'),
        ('uvx', 'Python包执行器', 'uvx example-tool'),
        ('本地', '本地可执行文件', './local-tool')
    ]
    
    for runtime, desc, cmd in runtime_data:
        row_cells = table.add_row().cells
        row_cells[0].text = runtime
        row_cells[1].text = desc  
        row_cells[2].text = cmd
    
    # 添加配置示例
    doc.add_heading('3. 配置示例', level=1)
    doc.add_paragraph('以下是典型的.env配置文件示例:')
    
    config_text = """
# AI模型配置
OPENAI_API_KEY=your_openai_key
OPENAI_BASE_URL=https://api.openai.com/v1
OPENAI_MODEL=gpt-4

# 数据库配置  
SUPABASE_URL=https://your-project.supabase.co
SUPABASE_SERVICE_ROLE_KEY=your_service_role_key
"""
    
    doc.add_paragraph(config_text)
    
    # 保存文档
    doc.save(TEST_DATA_DIR / "mcp_test_guide.docx")
    print("✅ Word文档生成完成")

def generate_powerpoint():
    """生成PowerPoint演示文稿"""
    print("🎤 生成PowerPoint演示文稿...")
    
    prs = Presentation()
    
    # 幻灯片1: 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "MCP测试框架"
    subtitle.text = "自动化Model Context Protocol工具测试解决方案\n" + datetime.now().strftime('%Y年%m月%d日')
    
    # 幻灯片2: 概览
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '框架特性'
    
    tf = body_shape.text_frame
    tf.text = '自动化部署'
    
    p = tf.add_paragraph()
    p.text = '智能测试生成'
    p.level = 0
    
    p = tf.add_paragraph()  
    p.text = '多格式报告'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '数据库集成'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '支持npx和uvx运行时'
    p.level = 0
    
    # 幻灯片3: 测试结果
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '测试结果统计'
    
    tf = body_shape.text_frame
    tf.text = '总测试数: 169个uvx工具'
    
    p = tf.add_paragraph()
    p.text = '成功率: 85%'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '平均测试时间: 15.3秒'
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = '发现25个高质量工具'
    p.level = 0
    
    # 保存演示文稿
    prs.save(TEST_DATA_DIR / "mcp_framework_presentation.pptx")
    print("✅ PowerPoint演示文稿生成完成")

def generate_images():
    """生成测试图片"""
    print("🖼️ 生成测试图片...")
    
    # 1. 创建简单的图表图片
    img = Image.new('RGB', (800, 600), color='white')
    draw = ImageDraw.Draw(img)
    
    # 绘制标题
    draw.text((50, 30), "MCP测试框架架构图", fill='black')
    
    # 绘制简单的框架图
    # 输入层
    draw.rectangle([50, 100, 200, 150], outline='blue', width=2)
    draw.text((60, 115), "GitHub URL", fill='blue')
    
    draw.rectangle([250, 100, 400, 150], outline='blue', width=2)
    draw.text((270, 115), "Package Name", fill='blue')
    
    # 处理层
    draw.rectangle([125, 200, 325, 250], outline='green', width=2)
    draw.text((180, 215), "MCP部署器", fill='green')
    
    # 测试层
    draw.rectangle([50, 300, 200, 350], outline='orange', width=2)
    draw.text((85, 315), "基础测试", fill='orange')
    
    draw.rectangle([250, 300, 400, 350], outline='orange', width=2)
    draw.text((285, 315), "AI智能测试", fill='orange')
    
    # 输出层
    draw.rectangle([125, 450, 325, 500], outline='red', width=2)
    draw.text((185, 465), "测试报告", fill='red')
    
    # 绘制连接线
    draw.line([125, 150, 175, 200], fill='black', width=2)
    draw.line([325, 150, 275, 200], fill='black', width=2)
    draw.line([225, 250, 125, 300], fill='black', width=2)
    draw.line([225, 250, 325, 300], fill='black', width=2)
    draw.line([125, 350, 175, 450], fill='black', width=2)
    draw.line([325, 350, 275, 450], fill='black', width=2)
    
    img.save(TEST_DATA_DIR / "framework_architecture.png")
    
    # 2. 创建一个简单的数据可视化图片
    chart_img = Image.new('RGB', (600, 400), color='white')
    chart_draw = ImageDraw.Draw(chart_img)
    
    chart_draw.text((200, 20), "MCP工具测试成功率", fill='black')
    
    # 绘制简单的柱状图
    categories = ['npx工具', 'uvx工具', '本地工具']
    values = [88, 85, 92]  # 成功率百分比
    colors = ['blue', 'green', 'red']
    
    bar_width = 80
    bar_spacing = 150
    base_y = 350
    
    for i, (cat, val, color) in enumerate(zip(categories, values, colors)):
        x = 50 + i * bar_spacing
        bar_height = val * 2.5  # 缩放高度
        
        # 绘制柱子
        chart_draw.rectangle([x, base_y - bar_height, x + bar_width, base_y], fill=color)
        
        # 绘制标签
        chart_draw.text((x + 10, base_y + 10), cat, fill='black')
        chart_draw.text((x + 30, base_y - bar_height - 20), f'{val}%', fill='black')
    
    chart_img.save(TEST_DATA_DIR / "test_success_rate_chart.png")
    
    print("✅ 测试图片生成完成")

def generate_sample_files():
    """生成示例业务文件"""
    print("📋 生成示例业务文件...")
    
    # 生成不同格式的员工数据以便测试不同的MCP工具
    
    # 1. 简单的员工CSV（用于Excel工具测试）
    simple_employee_data = [
        ["姓名", "部门", "薪资"],
        ["张三", "技术部", 15000],
        ["李四", "产品部", 18000],
        ["王五", "设计部", 12000]
    ]
    
    with open(TEST_DATA_DIR / "simple_employees.csv", "w", newline="", encoding="utf-8") as f:
        import csv
        writer = csv.writer(f)
        writer.writerows(simple_employee_data)
    
    # 2. 产品信息JSON（用于API工具测试）
    products = {
        "products": [
            {"id": 1, "name": "笔记本电脑", "price": 5999, "category": "电子产品", "stock": 50},
            {"id": 2, "name": "无线鼠标", "price": 199, "category": "配件", "stock": 200},
            {"id": 3, "name": "机械键盘", "price": 599, "category": "配件", "stock": 80},
            {"id": 4, "name": "显示器", "price": 1999, "category": "电子产品", "stock": 30}
        ],
        "meta": {
            "total_products": 4,
            "last_updated": datetime.now().isoformat()
        }
    }
    
    with open(TEST_DATA_DIR / "products.json", "w", encoding="utf-8") as f:
        json.dump(products, f, indent=2, ensure_ascii=False)
    
    # 3. 配置XML文件
    root = ET.Element("configuration")
    
    database = ET.SubElement(root, "database")
    ET.SubElement(database, "host").text = "localhost"
    ET.SubElement(database, "port").text = "5432"
    ET.SubElement(database, "name").text = "mcp_test"
    
    api = ET.SubElement(root, "api")
    ET.SubElement(api, "base_url").text = "https://api.example.com"
    ET.SubElement(api, "timeout").text = "30"
    
    tree = ET.ElementTree(root)
    tree.write(TEST_DATA_DIR / "application_config.xml", encoding="utf-8", xml_declaration=True)
    
    print("✅ 示例业务文件生成完成")

def main():
    """主函数"""
    print("🎯 开始生成完整的MCP测试数据集...")
    print(f"📂 目标目录: {TEST_DATA_DIR}")
    
    try:
        generate_comprehensive_excel()
        generate_comprehensive_pdf() 
        generate_word_document()
        generate_powerpoint()
        generate_images()
        generate_sample_files()
        
        print(f"\n✅ 完整测试数据集生成成功！")
        print(f"📁 所有文件保存在: {TEST_DATA_DIR}")
        
        # 统计生成的文件
        files = list(TEST_DATA_DIR.glob("*"))
        total_size = sum(f.stat().st_size for f in files if f.is_file())
        
        print(f"\n📊 生成统计:")
        print(f"   📄 文件总数: {len(files)}")
        print(f"   💾 总大小: {total_size / 1024:.1f} KB")
        
        # 按类型分类显示文件
        file_types = {}
        for file in files:
            if file.is_file():
                ext = file.suffix.lower()
                if ext not in file_types:
                    file_types[ext] = []
                file_types[ext].append(file.name)
        
        print(f"\n📋 按类型分类:")
        for ext, filenames in sorted(file_types.items()):
            print(f"   {ext or '无扩展名'}: {len(filenames)}个")
            for name in filenames:
                print(f"      📄 {name}")
        
        print(f"\n🚀 测试数据已就绪！现在可以测试各种MCP工具了：")
        print(f"   • Excel工具: 使用 financial_report.xlsx 和 employee_management.xlsx")
        print(f"   • PDF工具: 使用 mcp_test_report.pdf") 
        print(f"   • Word工具: 使用 mcp_test_guide.docx")
        print(f"   • PowerPoint工具: 使用 mcp_framework_presentation.pptx")
        print(f"   • 图片工具: 使用 *.png 文件")
        print(f"   • CSV工具: 使用各种 *.csv 文件")
        
    except Exception as e:
        print(f"❌ 生成过程中出现错误: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()
