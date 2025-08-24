#!/usr/bin/env python3
"""
æµ‹è¯•æ•°æ®ç”Ÿæˆå™¨

ä¸ºMCPå·¥å…·æµ‹è¯•ç”Ÿæˆå„ç§æ ¼å¼çš„æ ·æœ¬æ–‡ä»¶ï¼š
- CSV: åŒ…å«å‘˜å·¥ä¿¡æ¯ã€é”€å”®æ•°æ®ç­‰
- Excel: å¤šå·¥ä½œè¡¨è´¢åŠ¡æ•°æ®
- Wordæ–‡æ¡£: åŒ…å«æ–‡æœ¬ã€è¡¨æ ¼ã€å›¾ç‰‡è¯´æ˜
- PDF: æŠ¥å‘Šæ ¼å¼æ–‡æ¡£
- PowerPoint: æ¼”ç¤ºæ–‡ç¨¿

ä½œè€…: AI Assistant  
æ—¥æœŸ: 2025-08-24
"""

import csv
import json
import os
from pathlib import Path
from datetime import datetime, timedelta
import random

# æµ‹è¯•æ•°æ®ç›®å½•
TEST_DATA_DIR = Path(__file__).parent.parent / "test_data"
TEST_DATA_DIR.mkdir(exist_ok=True)

def generate_csv_files():
    """ç”ŸæˆCSVæµ‹è¯•æ–‡ä»¶"""
    print("ğŸ“Š ç”ŸæˆCSVæ–‡ä»¶...")
    
    # 1. å‘˜å·¥ä¿¡æ¯è¡¨
    employees_data = [
        ["ID", "å§“å", "éƒ¨é—¨", "èŒä½", "è–ªèµ„", "å…¥èŒæ—¥æœŸ", "é‚®ç®±"],
        [1, "å¼ ä¸‰", "æŠ€æœ¯éƒ¨", "é«˜çº§å·¥ç¨‹å¸ˆ", 15000, "2022-01-15", "zhangsan@company.com"],
        [2, "æå››", "äº§å“éƒ¨", "äº§å“ç»ç†", 18000, "2021-08-20", "lisi@company.com"],
        [3, "ç‹äº”", "è®¾è®¡éƒ¨", "UIè®¾è®¡å¸ˆ", 12000, "2022-03-10", "wangwu@company.com"],
        [4, "èµµå…­", "æŠ€æœ¯éƒ¨", "å‰ç«¯å·¥ç¨‹å¸ˆ", 13000, "2022-06-01", "zhaoliu@company.com"],
        [5, "é’±ä¸ƒ", "è¿è¥éƒ¨", "è¿è¥ä¸“å‘˜", 10000, "2022-09-15", "qianqi@company.com"],
        [6, "å­™å…«", "æŠ€æœ¯éƒ¨", "æ¶æ„å¸ˆ", 25000, "2020-05-10", "sunba@company.com"],
        [7, "å‘¨ä¹", "äººäº‹éƒ¨", "æ‹›è˜ç»ç†", 14000, "2021-12-01", "zhoujiu@company.com"],
        [8, "å´å", "è´¢åŠ¡éƒ¨", "ä¼šè®¡", 11000, "2022-02-28", "wushi@company.com"],
    ]
    
    with open(TEST_DATA_DIR / "employees.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(employees_data)
    
    # 2. é”€å”®æ•°æ®
    sales_data = [["æ—¥æœŸ", "äº§å“", "é”€é‡", "å•ä»·", "æ€»é¢", "é”€å”®å‘˜"]]
    products = ["ç¬”è®°æœ¬ç”µè„‘", "å°å¼æœº", "æ˜¾ç¤ºå™¨", "é”®ç›˜", "é¼ æ ‡", "è€³æœº"]
    salespeople = ["å°æ˜", "å°çº¢", "å°åˆš", "å°ç¾", "å°å¼º"]
    
    base_date = datetime(2024, 1, 1)
    for i in range(100):
        date = base_date + timedelta(days=random.randint(0, 365))
        product = random.choice(products)
        quantity = random.randint(1, 20)
        unit_price = random.randint(500, 8000)
        total = quantity * unit_price
        salesperson = random.choice(salespeople)
        sales_data.append([
            date.strftime("%Y-%m-%d"),
            product,
            quantity,
            unit_price,
            total,
            salesperson
        ])
    
    with open(TEST_DATA_DIR / "sales_data.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(sales_data)
    
    # 3. MCPå·¥å…·æµ‹è¯•æ•°æ®
    mcp_tools_data = [
        ["å·¥å…·åç§°", "ç±»å‹", "æè¿°", "æ”¯æŒè¿è¡Œæ—¶", "æ˜¯å¦éœ€è¦APIå¯†é’¥", "GitHubåœ°å€"],
        ["Excel MCP Server", "æ•°æ®å¤„ç†", "ç”¨äºæ“ä½œExcelæ–‡ä»¶çš„MCPæœåŠ¡å™¨", "uvx", "å¦", "https://github.com/haris-musa/excel-mcp-server"],
        ["Context7", "æ–‡æ¡£æ£€ç´¢", "è·å–æœ€æ–°çš„åº“æ–‡æ¡£", "npx", "å¦", "https://github.com/upstash/context7"],
        ["Blender MCP", "3Då»ºæ¨¡", "é€šè¿‡MCPæ§åˆ¶Blender", "uvx", "å¦", "https://github.com/ahujasid/blender-mcp"],
        ["ElevenLabs MCP", "è¯­éŸ³åˆæˆ", "æ–‡æœ¬è½¬è¯­éŸ³æœåŠ¡", "uvx", "æ˜¯", "https://github.com/elevenlabs/elevenlabs-mcp"],
    ]
    
    with open(TEST_DATA_DIR / "mcp_tools.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(mcp_tools_data)
    
    print("âœ… CSVæ–‡ä»¶ç”Ÿæˆå®Œæˆ")

def generate_json_files():
    """ç”ŸæˆJSONæµ‹è¯•æ–‡ä»¶"""
    print("ğŸ“„ ç”ŸæˆJSONæ–‡ä»¶...")
    
    # 1. é…ç½®æ–‡ä»¶
    config_data = {
        "app_name": "MCP Test Framework",
        "version": "1.0.0",
        "database": {
            "host": "localhost",
            "port": 5432,
            "name": "mcp_test",
            "ssl": True
        },
        "api_endpoints": {
            "base_url": "https://api.example.com",
            "timeout": 30,
            "retry_count": 3
        },
        "features": {
            "smart_testing": True,
            "database_export": True,
            "report_generation": True
        }
    }
    
    with open(TEST_DATA_DIR / "config.json", "w", encoding="utf-8") as f:
        json.dump(config_data, f, indent=2, ensure_ascii=False)
    
    # 2. æµ‹è¯•ç»“æœæ•°æ®
    test_results = {
        "test_session": {
            "id": "test_20240824_001",
            "timestamp": datetime.now().isoformat(),
            "total_tests": 25,
            "passed": 20,
            "failed": 5,
            "duration": 120.5
        },
        "test_cases": [
            {
                "name": "MCPè¿æ¥æµ‹è¯•",
                "status": "passed",
                "duration": 2.3,
                "message": "æˆåŠŸå»ºç«‹MCPè¿æ¥"
            },
            {
                "name": "å·¥å…·åˆ—è¡¨è·å–",
                "status": "passed", 
                "duration": 1.8,
                "message": "æˆåŠŸè·å–15ä¸ªå¯ç”¨å·¥å…·"
            },
            {
                "name": "Excelæ–‡ä»¶æ“ä½œ",
                "status": "failed",
                "duration": 5.2,
                "message": "æ–‡ä»¶è·¯å¾„é”™è¯¯",
                "error": "Invalid filename: test.xlsx, must be an absolute path"
            }
        ]
    }
    
    with open(TEST_DATA_DIR / "test_results.json", "w", encoding="utf-8") as f:
        json.dump(test_results, f, indent=2, ensure_ascii=False)
    
    print("âœ… JSONæ–‡ä»¶ç”Ÿæˆå®Œæˆ")

def generate_text_files():
    """ç”Ÿæˆæ–‡æœ¬æµ‹è¯•æ–‡ä»¶"""
    print("ğŸ“ ç”Ÿæˆæ–‡æœ¬æ–‡ä»¶...")
    
    # 1. READMEæ–‡ä»¶
    readme_content = """# MCPæµ‹è¯•æ¡†æ¶

è¿™æ˜¯ä¸€ä¸ªç”¨äºæµ‹è¯•Model Context Protocol (MCP)å·¥å…·çš„è‡ªåŠ¨åŒ–æ¡†æ¶ã€‚

## ä¸»è¦åŠŸèƒ½

- ğŸš€ è‡ªåŠ¨åŒ–MCPå·¥å…·éƒ¨ç½² (npx/uvx)
- ğŸ§ª æ™ºèƒ½æµ‹è¯•ç”¨ä¾‹ç”Ÿæˆ
- ğŸ“Š è¯¦ç»†çš„æµ‹è¯•æŠ¥å‘Š
- ğŸ—„ï¸ æ•°æ®åº“ç»“æœå­˜å‚¨
- ğŸ” GitHubä»“åº“è¯„ä¼°

## å¿«é€Ÿå¼€å§‹

```bash
# å®‰è£…ä¾èµ–
uv sync

# æµ‹è¯•å•ä¸ªå·¥å…·
uv run python -m src.main test-url "https://github.com/example/mcp-tool"

# æ‰¹é‡æµ‹è¯•
uv run python -m src.main batch-test --input data/test.csv
```

## æ”¯æŒçš„æ–‡ä»¶æ ¼å¼

- CSV: å‘˜å·¥æ•°æ®ã€é”€å”®è®°å½•
- Excel: å¤šå·¥ä½œè¡¨è´¢åŠ¡æ•°æ®  
- JSON: é…ç½®æ–‡ä»¶ã€APIå“åº”
- PDF: æŠ¥å‘Šæ–‡æ¡£
- PowerPoint: æ¼”ç¤ºæ–‡ç¨¿

## è´¡çŒ®æŒ‡å—

æ¬¢è¿æäº¤Issueå’ŒPull Requestï¼

## è®¸å¯è¯

MIT License
"""
    
    with open(TEST_DATA_DIR / "README.md", "w", encoding="utf-8") as f:
        f.write(readme_content)
    
    # 2. æ—¥å¿—æ–‡ä»¶æ ·æœ¬
    log_content = """2024-08-24 10:30:15 | INFO     | å¯åŠ¨MCPæµ‹è¯•æ¡†æ¶
2024-08-24 10:30:16 | INFO     | åŠ è½½é…ç½®æ–‡ä»¶: config.json
2024-08-24 10:30:17 | INFO     | è¿æ¥æ•°æ®åº“æˆåŠŸ
2024-08-24 10:30:18 | INFO     | å¼€å§‹æµ‹è¯•å·¥å…·: excel-mcp-server
2024-08-24 10:30:19 | DEBUG    | æ‰§è¡Œå‘½ä»¤: uvx excel-mcp-server stdio
2024-08-24 10:30:21 | INFO     | MCPè¿æ¥å»ºç«‹æˆåŠŸ
2024-08-24 10:30:22 | INFO     | å‘ç°25ä¸ªå¯ç”¨å·¥å…·
2024-08-24 10:30:23 | WARNING  | å·¥å…·éœ€è¦ç»å¯¹è·¯å¾„å‚æ•°
2024-08-24 10:30:25 | ERROR    | æµ‹è¯•ç”¨ä¾‹æ‰§è¡Œå¤±è´¥: è·¯å¾„æ ¼å¼é”™è¯¯
2024-08-24 10:30:26 | INFO     | ç”Ÿæˆæµ‹è¯•æŠ¥å‘Šå®Œæˆ
2024-08-24 10:30:27 | INFO     | æ¸…ç†èµ„æºå®Œæˆ
"""
    
    with open(TEST_DATA_DIR / "test.log", "w", encoding="utf-8") as f:
        f.write(log_content)
    
    print("âœ… æ–‡æœ¬æ–‡ä»¶ç”Ÿæˆå®Œæˆ")

def generate_advanced_files():
    """ç”ŸæˆExcelã€Wordç­‰é«˜çº§æ ¼å¼æ–‡ä»¶ï¼ˆå¦‚æœåº“å¯ç”¨ï¼‰"""
    print("ğŸ“ˆ å°è¯•ç”Ÿæˆé«˜çº§æ ¼å¼æ–‡ä»¶...")
    
    # å°è¯•ç”ŸæˆExcelæ–‡ä»¶
    try:
        import pandas as pd
        
        # åˆ›å»ºå¤šå·¥ä½œè¡¨Excelæ–‡ä»¶
        with pd.ExcelWriter(TEST_DATA_DIR / "financial_report.xlsx", engine='openpyxl') as writer:
            # æ”¶å…¥è¡¨
            revenue_data = pd.DataFrame({
                'æœˆä»½': ['1æœˆ', '2æœˆ', '3æœˆ', '4æœˆ', '5æœˆ', '6æœˆ'],
                'äº§å“Aæ”¶å…¥': [50000, 52000, 48000, 55000, 58000, 60000],
                'äº§å“Bæ”¶å…¥': [30000, 31000, 29000, 32000, 35000, 38000],
                'æ€»æ”¶å…¥': [80000, 83000, 77000, 87000, 93000, 98000]
            })
            revenue_data.to_sheet(writer, sheet_name='æ”¶å…¥æŠ¥è¡¨', index=False)
            
            # æˆæœ¬è¡¨  
            cost_data = pd.DataFrame({
                'æœˆä»½': ['1æœˆ', '2æœˆ', '3æœˆ', '4æœˆ', '5æœˆ', '6æœˆ'],
                'äººå·¥æˆæœ¬': [25000, 25000, 25000, 26000, 26000, 27000],
                'ææ–™æˆæœ¬': [15000, 16000, 14000, 17000, 18000, 19000],
                'è¿è¥æˆæœ¬': [8000, 8200, 7800, 8500, 8800, 9000],
                'æ€»æˆæœ¬': [48000, 49200, 46800, 51500, 52800, 55000]
            })
            cost_data.to_sheet(writer, sheet_name='æˆæœ¬æŠ¥è¡¨', index=False)
        
        print("âœ… Excelæ–‡ä»¶ç”Ÿæˆå®Œæˆ")
        
    except ImportError:
        print("âš ï¸ pandasæœªå®‰è£…ï¼Œè·³è¿‡Excelæ–‡ä»¶ç”Ÿæˆ")
    
    # ç”Ÿæˆç®€å•çš„XMLæ–‡ä»¶
    xml_content = '''<?xml version="1.0" encoding="UTF-8"?>
<mcp_tools>
    <tool>
        <name>Excel MCP Server</name>
        <type>data_processing</type>
        <runtime>uvx</runtime>
        <requires_api_key>false</requires_api_key>
        <capabilities>
            <capability>read_excel</capability>
            <capability>write_excel</capability>
            <capability>format_cells</capability>
            <capability>create_charts</capability>
        </capabilities>
    </tool>
    <tool>
        <name>Context7</name>
        <type>documentation</type>
        <runtime>npx</runtime>
        <requires_api_key>false</requires_api_key>
        <capabilities>
            <capability>resolve_library_id</capability>
            <capability>get_library_docs</capability>
        </capabilities>
    </tool>
</mcp_tools>'''
    
    with open(TEST_DATA_DIR / "mcp_tools.xml", "w", encoding="utf-8") as f:
        f.write(xml_content)
    
    print("âœ… XMLæ–‡ä»¶ç”Ÿæˆå®Œæˆ")

def generate_binary_placeholders():
    """ç”ŸæˆäºŒè¿›åˆ¶æ–‡ä»¶çš„å ä½ç¬¦è¯´æ˜"""
    print("ğŸ“‹ ç”ŸæˆäºŒè¿›åˆ¶æ–‡ä»¶è¯´æ˜...")
    
    binary_info = """# äºŒè¿›åˆ¶æ–‡ä»¶è¯´æ˜

ç”±äºPythonç¯å¢ƒé™åˆ¶ï¼Œä»¥ä¸‹äºŒè¿›åˆ¶æ ¼å¼æ–‡ä»¶éœ€è¦æ‰‹åŠ¨åˆ›å»ºæˆ–ä½¿ç”¨ä¸“é—¨å·¥å…·ç”Ÿæˆï¼š

## PDFæ–‡ä»¶
æ¨èä½¿ç”¨ä»¥ä¸‹Pythonåº“ç”Ÿæˆï¼š
- reportlab: åˆ›å»ºå¤æ‚PDFæŠ¥å‘Š
- fpdf2: è½»é‡çº§PDFç”Ÿæˆ
- weasyprint: HTMLè½¬PDF

ç¤ºä¾‹å‘½ä»¤ï¼š
```bash
pip install reportlab
python -c "
from reportlab.pdfgen import canvas
c = canvas.Canvas('test_data/sample_report.pdf')
c.drawString(100, 750, 'MCPæµ‹è¯•æ¡†æ¶æŠ¥å‘Š')
c.drawString(100, 730, 'è¿™æ˜¯ä¸€ä¸ªæµ‹è¯•PDFæ–‡ä»¶')
c.save()
"
```

## Wordæ–‡æ¡£ (.docx)
æ¨èä½¿ç”¨python-docxåº“ï¼š
```bash
pip install python-docx
python -c "
from docx import Document
doc = Document()
doc.add_heading('MCPæµ‹è¯•æ–‡æ¡£', 0)
doc.add_paragraph('è¿™æ˜¯ä¸€ä¸ªæµ‹è¯•Wordæ–‡æ¡£ï¼ŒåŒ…å«äº†MCPå·¥å…·çš„ç›¸å…³ä¿¡æ¯ã€‚')
doc.save('test_data/sample_document.docx')
"
```

## PowerPointæ¼”ç¤ºæ–‡ç¨¿ (.pptx)
æ¨èä½¿ç”¨python-pptxåº“ï¼š
```bash
pip install python-pptx
python -c "
from pptx import Presentation
prs = Presentation()
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = 'MCPæµ‹è¯•æ¡†æ¶'
title_slide.shapes.placeholders[1].text = 'è‡ªåŠ¨åŒ–MCPå·¥å…·æµ‹è¯•è§£å†³æ–¹æ¡ˆ'
prs.save('test_data/sample_presentation.pptx')
"
```

## å›¾ç‰‡æ–‡ä»¶
å¯ä»¥ä½¿ç”¨PILç”Ÿæˆæµ‹è¯•å›¾ç‰‡ï¼š
```bash
pip install Pillow
python -c "
from PIL import Image, ImageDraw, ImageFont
img = Image.new('RGB', (800, 600), color='white')
draw = ImageDraw.Draw(img)
draw.text((50, 50), 'MCP Test Image', fill='black')
img.save('test_data/sample_image.png')
"
```
"""
    
    with open(TEST_DATA_DIR / "BINARY_FILES_GUIDE.md", "w", encoding="utf-8") as f:
        f.write(binary_info)
    
    print("âœ… äºŒè¿›åˆ¶æ–‡ä»¶è¯´æ˜ç”Ÿæˆå®Œæˆ")

def main():
    """ä¸»å‡½æ•°"""
    print("ğŸ¯ å¼€å§‹ç”ŸæˆMCPæµ‹è¯•æ•°æ®æ–‡ä»¶...")
    print(f"ğŸ“‚ ç›®æ ‡ç›®å½•: {TEST_DATA_DIR}")
    
    generate_csv_files()
    generate_json_files() 
    generate_text_files()
    generate_advanced_files()
    generate_binary_placeholders()
    
    print(f"\nâœ… æµ‹è¯•æ•°æ®ç”Ÿæˆå®Œæˆï¼")
    print(f"ğŸ“ ç”Ÿæˆçš„æ–‡ä»¶ä¿å­˜åœ¨: {TEST_DATA_DIR}")
    
    # åˆ—å‡ºç”Ÿæˆçš„æ–‡ä»¶
    files = list(TEST_DATA_DIR.glob("*"))
    if files:
        print("\nğŸ“‹ ç”Ÿæˆçš„æ–‡ä»¶åˆ—è¡¨:")
        for file in sorted(files):
            size = file.stat().st_size if file.is_file() else 0
            print(f"   ğŸ“„ {file.name} ({size} bytes)")
    
    print(f"\nğŸš€ ç°åœ¨æ‚¨å¯ä»¥ä½¿ç”¨è¿™äº›æµ‹è¯•æ–‡ä»¶æ¥æµ‹è¯•å„ç§MCPå·¥å…·äº†ï¼")

if __name__ == "__main__":
    main()
